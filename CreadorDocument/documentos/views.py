from django.contrib.auth import logout
from django.shortcuts import render, get_object_or_404, redirect
from django.urls import reverse
from .models import Documento
import zipfile
from django.core.files import File
import tempfile
from io import BytesIO
from pptx import Presentation
import fitz  
from django.views.decorators.http import require_http_methods
from django.http import FileResponse
from io import BytesIO
from .forms import DocumentoForm
from pdf2image import convert_from_bytes
from django.core.files.storage import default_storage
from pdf2image import convert_from_path
from django.contrib import messages
from django.contrib.auth.decorators import login_required
import os
from google.oauth2.credentials import Credentials
from django.conf import settings
import json
from django.http import HttpResponse, HttpResponseRedirect
import csv
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from pptx import Presentation, slide
from googleapiclient.http import MediaFileUpload
from googleapiclient.discovery import build
import mimetypes
from uuid import uuid4

EXTENSION_MAP = {
    'word': 'docx',
    'txt': 'txt',
    'excel': 'xlsx',
    'cal': 'xls',
    'ppt': 'pptx',
    'pdf': 'pdf',
    'html': 'html',
    'php': 'php',
    'css': 'css',
    'js': 'js',
}

def obtener_nombre_archivo_con_extension(titulo, extension):
    EXTENSION_MAP = {
        'word': 'docx',
        'txt': 'txt',
        'excel': 'xlsx',
        'cal': 'xls',
        'ppt': 'pptx',
        'pdf': 'pdf',
        'html': 'html',
        'php': 'php',
        'css': 'css',
        'js': 'js',
    }

    ext_real = EXTENSION_MAP.get(extension.lower(), extension.lower())  # Usa extensión real o la misma si no está en el mapa
    nombre_archivo = f"{titulo}.{ext_real}"
    return nombre_archivo

@login_required
def lista_documentos(request):
    tipo = request.GET.get('tipo')  # word, excel, etc.
    documentos = Documento.objects.filter(usuario=request.user)
    tipos_presentes = set(documentos.values_list('extension', flat=True))
    if tipo:
        if tipo == "word":
            documentos = documentos.filter(extension__in=['word', 'txt'])
        else:
            documentos = documentos.filter(extension=tipo)
    return render(request, 'documentos/lista.html', {'documentos': documentos, 'tipos_presentes': tipos_presentes})

@login_required
def crear_documento(request):
    if request.method == 'POST':
        titulo = request.POST.get('titulo')
        extension = request.POST.get('extension')
        contenido = request.POST.get('contenido')
        usuario = request.user

        # Generamos el nombre del archivo y la ruta donde se guardará
        nombre_archivo = obtener_nombre_archivo_con_extension(titulo, extension)
        path_archivo = os.path.join(settings.MEDIA_ROOT, nombre_archivo)
        ruta_completa = os.path.join(settings.BASE_DIR, path_archivo)

        # Crear carpeta si no existe
        os.makedirs(os.path.dirname(ruta_completa), exist_ok=True)

        if extension == 'ppt':
            # Crear una nueva presentación de PowerPoint
            prs = Presentation()

            # Asumimos que cada diapositiva en el contenido está separada por dos saltos de línea
            diapositivas = contenido.split('\n\n')  # Cada "diapositiva" separada por doble salto de línea

            for diapositiva in diapositivas:
                slide = prs.slides.add_slide(prs.slide_layouts[1])  # Añadir una diapositiva de tipo título y contenido
                shapes = slide.shapes
                title = shapes.title
                content = shapes.placeholders[1]

                # Título de la diapositiva (puedes personalizarlo)
                title.text = f"Diapositiva {diapositivas.index(diapositiva) + 1}"
                content.text = diapositiva  # El contenido de la diapositiva

            # Guardar el archivo PowerPoint
            prs.save(ruta_completa)

        else:
            # Para otros tipos de documento, guardamos como texto
            with open(ruta_completa, 'w', encoding='utf-8') as f:
                f.write(contenido)

        # Crear el objeto Documento en la base de datos
        documento = Documento.objects.create(
            usuario=usuario,
            titulo=titulo,
            extension=extension,
            archivo_path=path_archivo
        )

        return redirect('editar_documento', doc_id=documento.id)

    return render(request, 'documentos/crear.html')
 

def csv_to_json(csv_data):
    reader = csv.reader(csv_data.splitlines())
    headers = next(reader)  # Asumiendo que la primera fila son los encabezados
    json_data = []

    for row in reader:
        row_data = {}
        for i, header in enumerate(headers):
            row_data[header] = row[i]
        json_data.append(row_data)

    return json_data
@login_required
def editar_documento(request, doc_id):
    documento = get_object_or_404(Documento, id=doc_id)

    if documento.extension in ['ppt', 'pptx']:
        prs = Presentation()
        prs.save(documento.archivo_path.path)
        return redirect('editar_presentacion', doc_id=documento.id)

    if request.method == 'POST':
        nuevo_contenido = request.POST.get('contenido')
        accion = request.POST.get('accion')  # 'guardar' o 'cerrar'

        try:
            if documento.extension in ['xlsx', 'excel', 'spreadsheet']:
                json_contenido = json.loads(nuevo_contenido)
                with open(documento.archivo_path.path, 'w', encoding='utf-8') as f:
                    json.dump(json_contenido, f, ensure_ascii=False, indent=4)
            else:
                with open(documento.archivo_path.path, 'w', encoding='utf-8') as f:
                    f.write(nuevo_contenido)
        except Exception as e:
            return render(request, 'documentos/error.html', {'mensaje': f'Error al guardar: {e}'})

        if accion == 'cerrar':
            return redirect('lista_documentos')
        else:
            messages.success(request, "Documento guardado.")
            return redirect('editar_documento', doc_id=documento.id)

    # Mostrar contenido del archivo
    if documento.extension in ['html', 'css', 'JavaScript', 'php']:
        with open(documento.archivo_path.path, 'r', encoding='utf-8') as f:
            contenido = f.read()
        return render(request, 'documentos/editar_codigo.html', {'documento': documento, 'contenido': contenido})
    
    elif documento.extension in ['xlsx', 'excel', 'spreadsheet']:
        try:
            with open(documento.archivo_path.path, 'r', encoding='utf-8') as f:
                contenido = f.read().strip()

            if contenido:
                try:
                    contenido_json = json.loads(contenido)
                    # Validar estructura mínima para evitar errores en LuckySheet
                    if not isinstance(contenido_json, list) or not all('celldata' in hoja for hoja in contenido_json):
                        raise ValueError("Estructura no válida")
                except Exception:
                    # Si falla, forzar una estructura básica por defecto
                    contenido_json = [{
                        "name": "Hoja1",
                        "color": "",
                        "index": 0,
                        "status": 1,
                        "order": 0,
                        "celldata": [],
                        "row": 30,
                        "column": 15
                    }]
            else:
                contenido_json = [{
                    "name": "Hoja1",
                    "color": "",
                    "index": 0,
                    "status": 1,
                    "order": 0,
                    "celldata": [],
                    "row": 30,
                    "column": 15
                }]
        except Exception as e:
            print(f"Error al leer el archivo: {e}")
            contenido_json = [{
                "name": "Hoja1",
                "color": "",
                "index": 0,
                "status": 1,
                "order": 0,
                "celldata": [],
                "row": 30,
                "column": 15
            }]

        return render(request, 'documentos/editar_excel.html', {
            'documento': documento,
            'contenido': json.dumps(contenido_json)
        })

    else:
        with open(documento.archivo_path.path, 'r', encoding='utf-8') as f:
            contenido = f.read()
        return render(request, 'documentos/editar_wysiwyg.html', {'documento': documento, 'contenido': contenido})
@login_required
@csrf_exempt  # Asegúrate de manejar el CSRF token correctamente en producción
def guardar_hoja_calculo(request, doc_id):
    documento = get_object_or_404(Documento, id=doc_id)
    
    if request.method == 'POST':
        try:
            # Recibimos los datos JSON de la solicitud
            datos_json = json.loads(request.body)
            contenido = datos_json.get('contenido')
            
            if contenido:
                # Guardar el contenido en el archivo correspondiente
                with open(documento.archivo_path.path, 'w', encoding='utf-8') as f:
                    f.write(contenido)  # Escribe el contenido nuevo en el archivo
                with open(documento.archivo_path.path, 'r', encoding='utf-8') as f:
                    verificacion = f.read()
                    print("Contenido después de guardar:", verificacion)

                return JsonResponse({"status": "success", "message": "Guardado exitosamente"})
            else:
                return JsonResponse({"status": "error", "message": "No se proporcionó contenido"})
        except Exception as e:
            return JsonResponse({"status": "error", "message": f"Error al guardar: {str(e)}"})
    
    return JsonResponse({"status": "error", "message": "Método no permitido"})

def editar_presentacion(request, doc_id):
    doc = get_object_or_404(Documento, id=doc_id)

    try:
        with doc.archivo_path.open('rb') as f:
            contenido = f.read()

        prs = Presentation(BytesIO(contenido))
    except Exception:
        return render(request, "documentos/error.html", {
            "mensaje": "El archivo de presentación no se pudo abrir correctamente."
        })

    if len(prs.slides) == 0:
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = "Primera diapositiva"
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = "Contenido de ejemplo"

        output = BytesIO()
        prs.save(output)
        output.seek(0)

        doc.archivo_path.save(os.path.basename(doc.archivo_path.name), File(output), save=True)

    slides_data = []
    for slide_index, slide in enumerate(prs.slides):
        contenido = []
        for idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                contenido.append({
                    "idx": idx,
                    "text": shape.text
                })
        slides_data.append({
            "index": slide_index,
            "contenido": contenido
        })

    return render(request, "documentos/editar_presentacion.html", {
        "documento": doc,
        "slides": slides_data,
    })


def guardar_diapositiva(request, doc_id, slide_index):
    documento = get_object_or_404(Documento, id=doc_id)

    if request.method == "POST":
        try:
            with documento.archivo_path.open('rb') as f:
                contenido = f.read()

            prs = Presentation(BytesIO(contenido))

            if slide_index >= len(prs.slides):
                return render(request, "documentos/error.html", {
                    "mensaje": "La diapositiva que intentas editar no existe."
                })

            slide = prs.slides[slide_index]
            for idx, shape in enumerate(slide.shapes):
                if shape.has_text_frame:
                    key = f'texto_{slide_index}_{idx}'
                    if key in request.POST:
                        shape.text = request.POST[key]

            output = BytesIO()
            prs.save(output)
            output.seek(0)

            documento.archivo_path.save(os.path.basename(documento.archivo_path.name), File(output), save=True)

        except Exception as e:
            return render(request, "documentos/error.html", {
                "mensaje": f"No se pudo abrir o guardar el archivo PPT: {e}"
            })

    return redirect('editar_presentacion', doc_id=doc_id)


def eliminar_diapositiva(request, doc_id, slide_index):
    documento = get_object_or_404(Documento, id=doc_id)

    if request.method == "POST":
        try:
            with documento.archivo_path.open('rb') as f:
                contenido = f.read()

            prs = Presentation(BytesIO(contenido))

            slides = prs.slides._sldIdLst
            if slide_index < 0 or slide_index >= len(slides):
                return render(request, "documentos/error.html", {
                    "mensaje": "La diapositiva que intentas eliminar no existe."
                })

            slide_id = slides[slide_index]
            slides.remove(slide_id)

            output = BytesIO()
            prs.save(output)
            output.seek(0)

            documento.archivo_path.save(os.path.basename(documento.archivo_path.name), File(output), save=True)

        except Exception as e:
            return render(request, "documentos/error.html", {
                "mensaje": f"No se pudo abrir o guardar el archivo PPT: {e}"
            })

    return redirect('editar_presentacion', doc_id=doc_id)


def exportar_a_drive(request, doc_id):
    documento = get_object_or_404(Documento, id=doc_id)
    archivo_path = documento.archivo_path.path

    creds_json = request.session.get('google_credentials')
    if not creds_json:
        return HttpResponse("No hay credenciales en sesión.")

    creds_data = json.loads(creds_json)

    creds = Credentials(
        token=creds_data['token'],
        refresh_token=creds_data.get('refresh_token'),
        token_uri=creds_data['token_uri'],
        client_id=creds_data['client_id'],
        client_secret=creds_data['client_secret'],
        scopes=creds_data['scopes'],
    )

    service = build('drive', 'v3', credentials=creds)

    ext_real = EXTENSION_MAP.get(documento.extension, 'bin')
    nombre_archivo = f"{documento.titulo}.{ext_real}"

    mime_type, _ = mimetypes.guess_type(nombre_archivo)

    file_metadata = {
        'name': nombre_archivo,
        'mimeType': mime_type or 'application/octet-stream',
    }

    media = MediaFileUpload(archivo_path, mimetype=mime_type, resumable=True)

    service.files().create(
        body=file_metadata,
        media_body=media,
        fields='id'
    ).execute()

    return redirect('lista_documentos')


def obtener_servicio_drive(request):
    creds_json = request.session.get('google_credentials')
    if not creds_json:
        print("No hay credenciales almacenadas en la sesión.")
        return None

    try:
        creds_data = json.loads(creds_json)
        creds = Credentials(
            token=creds_data.get('token'),
            refresh_token=creds_data.get('refresh_token'),
            token_uri=creds_data.get('token_uri'),
            client_id=creds_data.get('client_id'),
            client_secret=creds_data.get('client_secret'),
            scopes=creds_data.get('scopes'),
        )
        servicio = build('drive', 'v3', credentials=creds)
        return servicio
    except Exception as e:
        print(f"Error al crear el servicio de Drive: {e}")
        return None

def borrar_archivo_drive(request, nombre_archivo_con_extension):
    servicio = obtener_servicio_drive(request)
    if not servicio:
        print("No se pudo obtener servicio de Drive para borrar archivo.")
        return False

    try:
        resultados = servicio.files().list(
            q=f"name = '{nombre_archivo_con_extension}' and trashed = false",
            spaces='drive',
            fields="files(id, name)"
        ).execute()
        archivos = resultados.get('files', [])
    except Exception as e:
        print(f"Error buscando archivos en Drive: {e}")
        return False

    if not archivos:
        print(f"No se encontró ningún archivo con el nombre: {nombre_archivo_con_extension}")
        return False

    for archivo in archivos:
        try:
            servicio.files().delete(fileId=archivo['id']).execute()
        except Exception as e:
            return False

    return True

def borrar_documento(request, doc_id):
    documento = get_object_or_404(Documento, id=doc_id)

    if request.method == 'POST':
        # Borra archivo local si existe
        archivo_path = documento.archivo_path.path
        if os.path.exists(archivo_path):
            try:
                os.remove(archivo_path)
                print(f"Archivo local eliminado: {archivo_path}")
            except Exception as e:
                messages.error(request, f"No se pudo borrar el archivo local: {e}")
                return redirect('lista_documentos')
        else:
            print("Archivo local no encontrado, continúa con borrado en Drive si aplica.")

        # Intenta borrar en Google Drive
        nombre_archivo_con_extension = obtener_nombre_archivo_con_extension(documento.titulo, documento.extension)
        exito_drive = borrar_archivo_drive(request, nombre_archivo_con_extension)

        if exito_drive:
            print("Archivo borrado en Google Drive correctamente.")
        else:
            print("No se pudo borrar el archivo en Google Drive o no existía.")

        # Finalmente borra el registro en la base de datos
        documento.delete()
        messages.success(request, "Documento borrado correctamente (local y Drive si aplicaba).")
        return redirect('lista_documentos')
# nuevo
def cerrar_sesion(request):
    # Elimina credenciales OAuth de la sesión si existen
    if 'google_credentials' in request.session:
        del request.session['google_credentials']

    # Cierra la sesión de Django
    logout(request)
    return redirect('inicio')
@login_required
def configurar_sesion(request):
    usuario_email = request.user.email if request.user.is_authenticated else ''
    return render(request, 'app/configurar_sesion.html', {'usuario_email': usuario_email})
@login_required
def importar_documento(request):
    if request.method == 'POST':
        archivo = request.FILES.get('archivo')
        if archivo:
            usuario = request.user
            nombre_archivo = archivo.name
            titulo, extension = os.path.splitext(nombre_archivo)
            extension = extension[1:].lower()  # quitar el punto y pasar a minúsculas

            # Nombre y ruta donde guardar el archivo en MEDIA_ROOT
            nombre_archivo_guardar = obtener_nombre_archivo_con_extension(titulo, extension)
            ruta_guardar = os.path.join(settings.MEDIA_ROOT, nombre_archivo_guardar)

            # Crear carpeta si no existe
            os.makedirs(os.path.dirname(ruta_guardar), exist_ok=True)

            # Guardar archivo subido en la ruta correspondiente
            with open(ruta_guardar, 'wb+') as destino:
                for chunk in archivo.chunks():
                    destino.write(chunk)

            # Crear objeto Documento en la base de datos
            Documento.objects.create(
                usuario=usuario,
                titulo=titulo,
                extension=extension,
                archivo_path=nombre_archivo_guardar,
            )

            messages.success(request, f'Documento "{nombre_archivo}" importado correctamente.')
        else:
            messages.error(request, 'No se recibió ningún archivo.')

    return redirect('lista_documentos')
def ordenar_documentos(request):
    ordenar_por = request.GET.get('ordenar_por', 'nombre')
    documentos = Documento.objects.filter(usuario=request.user)
    tipo = request.GET.get('tipo')  # word, excel, etc.
    documentos = Documento.objects.filter(usuario=request.user)
    if tipo:
        if tipo == "word":
            documentos = documentos.filter(extension__in=['word', 'txt'])
        else:
            documentos = documentos.filter(extension=tipo)
    if ordenar_por == 'nombre':
        documentos = documentos.order_by('titulo')
    return render(request, 'documentos/lista.html', {'documentos': documentos})

def pdf_a_imagen(request):
    if request.method == 'POST':
        archivo_pdf = request.FILES['archivo']
        doc = fitz.open(stream=archivo_pdf.read(), filetype="pdf")

        # Crear un buffer para el zip en memoria
        zip_buffer = BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w') as zip_file:
            for i in range(doc.page_count):
                pagina = doc.load_page(i)
                imagen = pagina.get_pixmap()

                # Obtener bytes PNG de la imagen
                img_bytes = imagen.tobytes("png")

                # Nombre del archivo para cada página
                nombre_img = f"pagina_{i+1}.png"

                # Agregar la imagen al zip
                zip_file.writestr(nombre_img, img_bytes)

        zip_buffer.seek(0)

        # Devolver el zip para descargar
        return FileResponse(zip_buffer, as_attachment=True, filename='imagenes_pdf.zip')

    return render(request, 'documentos/pdf_a_imagen.html')
def agregar_diapositiva(request, doc_id):
    documento = get_object_or_404(Documento, id=doc_id)

    try:
        with documento.archivo_path.open('rb') as f:
            contenido = f.read()

        prs = Presentation(BytesIO(contenido))

        slide_layout = prs.slide_layouts[1]
        nueva_slide = prs.slides.add_slide(slide_layout)

        if nueva_slide.shapes.title:
            nueva_slide.shapes.title.text = "Nueva diapositiva"
        if len(nueva_slide.shapes.placeholders) > 1:
            nueva_slide.shapes.placeholders[1].text = "Contenido de la nueva diapositiva"

        output = BytesIO()
        prs.save(output)
        output.seek(0)

        documento.archivo_path.save(os.path.basename(documento.archivo_path.name), File(output), save=True)

    except Exception as e:
        return render(request, "documentos/error.html", {
            "mensaje": f"No se pudo abrir o guardar el archivo PPT: {e}"
        })

    return redirect('editar_presentacion', doc_id=doc_id)
def guardar_y_volver_lista(request, doc_id):
    documento = get_object_or_404(Documento, id=doc_id)

    if request.method == "POST":
        try:
            with documento.archivo_path.open('rb') as f:
                contenido = f.read()

            prs = Presentation(BytesIO(contenido))

            for slide_index, slide in enumerate(prs.slides):
                for idx, shape in enumerate(slide.shapes):
                    if shape.has_text_frame:
                        key = f'texto_{slide_index}_{idx}'
                        if key in request.POST:
                            shape.text = request.POST[key]

            output = BytesIO()
            prs.save(output)
            output.seek(0)

            # Guardar con File() de Django
            documento.archivo_path.save(os.path.basename(documento.archivo_path.name), File(output), save=True)

            print("Archivo guardado correctamente.")

        except Exception as e:
            print("Error guardando presentación:", e)
            return render(request, "documentos/error.html", {
                "mensaje": f"No se pudo guardar el archivo PPT: {e}"
            })

    return redirect('lista_documentos')
def guardar_contenido_documento(request, doc_id):
    documento = get_object_or_404(Documento, id=doc_id)

    if request.method == 'POST':
        if request.content_type != 'application/json':
            return JsonResponse({"status": "error", "message": "Tipo de contenido no válido"})

        try:
            datos_json = json.loads(request.body)
            contenido = datos_json.get('contenido')

            if isinstance(contenido, str) and contenido.strip():
                with open(documento.archivo_path.path, 'w', encoding='utf-8') as f:
                    f.write(contenido)

                return JsonResponse({"status": "success", "message": "Guardado correctamente"})
            else:
                return JsonResponse({"status": "error", "message": "Contenido vacío o no válido"})
        except Exception as e:
            return JsonResponse({"status": "error", "message": f"Error al guardar: {str(e)}"})

    return JsonResponse({"status": "error", "message": "Método no permitido"})